from pptx import Presentation
import os

PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v2.7.pptx"

if os.path.exists(PATH):
    prs = Presentation(PATH)
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  - Text: {shape.text[:50]}...")
else:
    print("File not found")
