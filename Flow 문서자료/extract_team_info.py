from pptx import Presentation
import os

PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"

if os.path.exists(PATH):
    prs = Presentation(PATH)
    target_indices = [16, 20] # Slide 17 and Slide 21
    for idx in target_indices:
        if idx < len(prs.slides):
            slide = prs.slides[idx]
            print(f"\n--- Slide {idx+1} (Index {idx}) ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"Shape: {shape.name}")
                    print(f"Text: {shape.text}")
                    print("-" * 20)
else:
    print(f"File not found: {PATH}")
