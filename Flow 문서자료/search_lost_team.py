from pptx import Presentation
import os

FILES = [
    r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.5.pptx",
    r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.4_Executive.pptx"
]

for PATH in FILES:
    if os.path.exists(PATH):
        prs = Presentation(PATH)
        print(f"\n--- Checking File: {os.path.basename(PATH)} ---")
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and ("Member" in shape.text or "Team" in shape.text or "팀" in shape.text):
                    print(f"Slide {i+1}: {shape.text[:100]}...")
