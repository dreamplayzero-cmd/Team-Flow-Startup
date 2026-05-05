from pptx import Presentation
import os

PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"

if os.path.exists(PATH):
    prs = Presentation(PATH)
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} (Index {i}) ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  Shape [{shape.name}]: {shape.text}")
else:
    print(f"File not found: {PATH}")
