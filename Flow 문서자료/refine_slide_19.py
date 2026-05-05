from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
DEMO_VIDEO_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\모바일 시연영상.mp4"
IMG_MOBILE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5\mobile_gis_mockup_v3_4_1777363524637.png"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

def set_green_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    # Dark green for premium look
    fill.fore_color.rgb = RGBColor(10, 45, 30) 

def add_blue_link(slide, text, left, top, width, height, link_path):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 255) # Pure Blue
    run.font.underline = True
    
    # Add hyperlink to the whole textbox or run
    video_url = f"file:///{link_path.replace(chr(92), '/')}"
    run.hyperlink.address = video_url
    return txBox

# --- Targeted Update (Slide 19 only) ---
# Slide 19 is index 18
if len(prs.slides) > 18:
    slide = prs.slides[18]
    
    # 1. Clear Slide 19 completely
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)
    
    # 2. Set Theme
    set_green_theme(slide)
    
    # 3. Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE FACE: SERVICE DEMONSTRATION"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(222, 255, 154) # Gold/Lime accent
    
    # 4. Add Full Mobile Mockup
    if os.path.exists(IMG_MOBILE):
        # Position in center-right
        slide.shapes.add_picture(IMG_MOBILE, Inches(5), Inches(1), height=Inches(6))
    
    # 5. Add "시연영상 재생" Blue Link (Simple and Direct)
    add_blue_link(slide, "시연영상 재생", Inches(0.5), Inches(3.5), Inches(4), Inches(1), DEMO_VIDEO_PATH)
    
    # Add some descriptive text for balance
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = "현장의 감각을 데이터로\n전환하는 직관적 UX 시연"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(245, 245, 245)

# Save
prs.save(PATH)
print(f"Presentation v3.6 Refined (Slide 19) saved to: {PATH}")
