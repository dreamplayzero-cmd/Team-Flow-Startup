from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
PRESENTATION_NAME = "Flow_Final_Project_Report_v2.7_Premium.pptx"
BASE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow"
DOCS_PATH = os.path.join(BASE_PATH, "Flow 문서자료")
IMG_PATH = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
DNA_PATH = os.path.join(BASE_PATH, r"backend\assets\visual_dna")

# --- Initialize Presentation ---
prs = Presentation()

# Helper for dark theme
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_text(slide, text, size=44, color=(222, 255, 154)): # Sovereign Gold
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(size)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# 1. Slide: Title
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)

title = slide.shapes.title
title.text = "THE SOVEREIGN INSIGHT"
title.text_frame.paragraphs[0].font.size = Pt(60)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)

subtitle = slide.placeholders[1]
subtitle.text = "Strategic Presentation Deck v2.7\nYouth Startup Flow Analysis Engine"
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(245, 245, 245)

# Add Cover Image
cover_img = os.path.join(IMG_PATH, "strategic_presentation_cover_1777346028339.png")
if os.path.exists(cover_img):
    slide.shapes.add_picture(cover_img, Inches(0.5), Inches(4.5), height=Inches(2.5))

# 2. Slide: Multi-Platform Ecosystem
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)
add_title_text(slide, "MULTI-PLATFORM ECOSYSTEM")

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Web Dashboard (React + Vite)"
p = tf.add_paragraph()
p.text = "- High-Resolution Geo-Data Visualization\n- Dynamic Scenario Simulation"
p.level = 1

p = tf.add_paragraph()
p.text = "Mobile Insight App (Flutter)"
p = tf.add_paragraph()
p.text = "- Real-time Scouting & AR Support\n- Persona-based Personalized Analysis"
p.level = 1

# 3. Slide: Scoring Engine
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)
add_title_text(slide, "INTEGRATED SCORING ENGINE LOGIC")

score_img = os.path.join(IMG_PATH, "score_summation_logic_viz_1777346419731.png")
if os.path.exists(score_img):
    slide.shapes.add_picture(score_img, Inches(5), Inches(1.5), height=Inches(4))

content = slide.placeholders[1]
content.text = "4 Core Pillars:\n1. Population (실시간 유동인구)\n2. Demand (네이버 검색 수요)\n3. Trend (블로그 감성 분석)\n4. Competition (상가 생존 시뮬레이션)\n\nFinal Score = Public(60%) + SNS(40%)"

# 4. Slide: District Precision Report
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)
add_title_text(slide, "DISTRICT PRECISION REPORT: SEONGSU")

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "AI Match Score: 84.5 pts"
p = tf.add_paragraph()
p.text = "Gentrification Risk: LOW"
p = tf.add_paragraph()
p.text = "Est. BEP Period: 18 months"
p = tf.add_paragraph()
p.text = "\nStrategic Insight:\n'성수동 연무장길 상권은 현재 감성 폭발형 성장 단계의 초입에 위치해 있습니다.'"

# 5. Slide: Visual DNA
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)
add_title_text(slide, "VISUAL DNA RECOMMENDATIONS")

# Add some DNA images
dna_img1 = os.path.join(DNA_PATH, "SS_ID_02_Seongsu_Cafe_IndustrialVintage.jpg")
if os.path.exists(dna_img1):
    slide.shapes.add_picture(dna_img1, Inches(0.5), Inches(2), height=Inches(3))

dna_img2 = os.path.join(DNA_PATH, "HN_MC_01_Hannam_Dining_ModernChic.jpg")
if os.path.exists(dna_img2):
    slide.shapes.add_picture(dna_img2, Inches(5), Inches(2), height=Inches(3))

# 6. Slide: Clustering & Future
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)
add_title_text(slide, "AI CLUSTERING & FUTURE SCENARIO")

cluster_img = os.path.join(IMG_PATH, "clustering_visualization_premium_1777346047876.png")
if os.path.exists(cluster_img):
    slide.shapes.add_picture(cluster_img, Inches(0.5), Inches(1.5), width=Inches(9))

# Save
output_path = os.path.join(DOCS_PATH, PRESENTATION_NAME)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
