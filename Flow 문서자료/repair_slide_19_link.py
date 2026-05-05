from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
DEMO_VIDEO_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\모바일 시연영상.mp4"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

# --- Targeted Repair (Slide 19 only) ---
# Slide 19 is index 18
if len(prs.slides) > 18:
    slide = prs.slides[18]
    
    # 1. Find the problematic link text
    # It might be in a shape with "시연영상" text
    for shape in slide.shapes:
        if hasattr(shape, "text") and "시연영상" in shape.text:
            # Found the shape, now fix all runs in the first paragraph
            tf = shape.text_frame
            tf.clear() # Clear existing text/formatting
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text = "시연영상 재생 (Click Here)"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 255) # Restoration to Blue
            run.font.underline = True
            # font name restoration (using Noto Sans KR or similar premium font)
            run.font.name = "Noto Sans KR" 
            
            # 2. Fix the Hyperlink
            video_url = f"file:///{DEMO_VIDEO_PATH.replace(chr(92), '/')}"
            run.hyperlink.address = video_url

# Save
prs.save(PATH)
print(f"Presentation v3.6 Repaired (Slide 19 Link) saved to: {PATH}")
