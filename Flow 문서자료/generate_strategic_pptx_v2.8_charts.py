from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
import os

# --- Configuration ---
PRESENTATION_NAME = "Flow_Final_Project_Report_v2.8_Strategic_Analysis.pptx"
BASE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow"
DOCS_PATH = os.path.join(BASE_PATH, "Flow 문서자료")
IMG_PATH = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"

# --- Initialize Presentation ---
prs = Presentation()

# Helper for dark theme
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_text(slide, text, size=36, color=(222, 255, 154)):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(size)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# 1. Slide: Cover
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_dark_background(slide)
title = slide.shapes.title
title.text = "THE SOVEREIGN INSIGHT v2.8"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)

# 2. Slide: Native PPT Chart (The "River Crossing" Technique)
slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
set_dark_background(slide)
add_title_text(slide, "DISTRICT COMPARATIVE ANALYSIS (NATIVE CHART)")

# Prepare Chart Data
chart_data = CategoryChartData()
chart_data.categories = ['Population', 'Demand', 'Trend', 'Competition']
chart_data.add_series('Seongsu-dong', (82.1, 75.4, 91.2, 45.8))
chart_data.add_series('Hannam-dong', (65.4, 88.2, 78.5, 62.1))

# Add Chart
x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)
chart = graphic_frame.chart

# Style Chart (Optional: Dark Theme adjustments)
# Note: Full chart styling via code is complex, but the data is now editable in PPT.
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# 3. Slide: Data Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_dark_background(slide)
add_title_text(slide, "STRATEGIC METRICS SUMMARY TABLE")

# Define Table Dimensions
rows, cols = 4, 4
left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(3)

# Add Table
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set Column Widths
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(2)
table.columns[2].width = Inches(2)
table.columns[3].width = Inches(2)

# Set Header Content
headers = ['District', 'Final Score', 'Risk Level', 'Opportunity']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(50, 50, 50)
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(222, 255, 154)

# Fill Data
data = [
    ['Seongsu-dong', '84.5', 'LOW', 'MZ Hotspot'],
    ['Hannam-dong', '78.2', 'MEDIUM', 'Luxury Dining'],
    ['Dosan Park', '92.0', 'LOW', 'High-End Retail']
]

for row_idx, row_data in enumerate(data):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)

# Save
output_path = os.path.join(DOCS_PATH, PRESENTATION_NAME)
prs.save(output_path)
print(f"Presentation v2.8 saved to: {output_path}")
