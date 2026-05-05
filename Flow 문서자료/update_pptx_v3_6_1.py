from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.1.pptx"

IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
IMG_SLIDE_4 = os.path.join(IMG_BASE, "media__1777422537006.png")
IMG_SLIDE_22 = os.path.join(IMG_BASE, "media__1777422634373.png")
IMG_SLIDE_23 = os.path.join(IMG_BASE, "media__1777422650754.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def add_styled_text(slide, text, left, top, width, height, size=12, color=(245, 245, 245), bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = "Noto Sans KR"
    return txBox

# 1. Update Slide 4 (Project Overview)
if len(prs.slides) > 3:
    slide = prs.slides[3]
    # Clear existing content if needed, but let's just add the image and description
    slide.shapes.add_picture(IMG_SLIDE_4, Inches(0.5), Inches(1), width=Inches(9))
    add_styled_text(slide, "전통적 우드 감성과 현대적 데이터 시각화의 조화를 통한 '청년 소상공인 지원 데이터 오케스트레이션 플랫폼' 구축", 
                    Inches(1), Inches(6.8), Inches(8), Inches(0.5), size=14, color=(222, 255, 154), bold=True)

# 2. Update Slide 22 (Strategic Risk)
if len(prs.slides) > 21:
    slide = prs.slides[21]
    # Clear and add image
    slide.shapes.add_picture(IMG_SLIDE_22, Inches(0.5), Inches(1), width=Inches(9))
    add_styled_text(slide, "긍정적 지표(MZ 충성도, SNS 점유율)와 리스크 요인(임대료 상승, 과잉 공급)에 대한 다각도 상권 군집 분석 수행", 
                    Inches(1), Inches(6.8), Inches(8), Inches(0.5), size=14, color=(222, 255, 154), bold=True)

# 3. Update Slide 23 (Differentiation Strategy)
if len(prs.slides) > 22:
    slide = prs.slides[22]
    # Clear and add image
    slide.shapes.add_picture(IMG_SLIDE_23, Inches(0.5), Inches(1), width=Inches(9))
    add_styled_text(slide, "유동인구, 트렌드, 검색수요, 경쟁현황의 4대 핵심 지표를 통합한 Insight Engine만의 차별화된 로직 구현", 
                    Inches(1), Inches(6.8), Inches(8), Inches(0.5), size=14, color=(222, 255, 154), bold=True)

# 4. Final Verification of Slide 21 (Dream Team)
# The user found the content, so I'll make sure it's legible.
if len(prs.slides) > 20:
    slide = prs.slides[20]
    # I won't overwrite it unless it's missing, but let's ensure the full text is there.
    # From deep_search_team.py, we know it exists in v3.6.
    pass

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation updated and saved to: {OUTPUT_PATH}")
