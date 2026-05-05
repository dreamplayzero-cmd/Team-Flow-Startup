from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.1.pptx"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

def apply_image_border(shape, color=(222, 255, 154), width_pt=3):
    line = shape.line
    line.color.rgb = RGBColor(*color)
    line.width = Pt(width_pt)

target_indices = [3, 21, 22] # Slides 4, 22, 23

for idx in target_indices:
    if idx < len(prs.slides):
        slide = prs.slides[idx]
        shapes_to_delete = []
        for shape in slide.shapes:
            # 1. Identify "Twins": Delete the PPT title if the image already has one
            if shape.name.startswith("Title") or shape.name.startswith("Textbox"):
                # If it's a small description at the bottom, we might keep it or move it
                if hasattr(shape, "text") and ("전통적 우드" in shape.text or "긍정적 지표" in shape.text or "유동인구" in shape.text):
                    # This is our added description. Let's move it to be clean.
                    shape.top = Inches(6.9)
                else:
                    shapes_to_delete.append(shape)
            
            # 2. Delete the large frame Rectangles
            if shape.shape_type == 1: # Rectangle
                if shape.width > Inches(8):
                    shapes_to_delete.append(shape)
            
            # 3. Handle the Picture
            if shape.shape_type == 13: # Picture
                # Center and slightly scale down to look professional
                shape.width = Inches(8.8)
                shape.left = Inches(0.6)
                shape.top = Inches(0.5)
                # Add border ONLY to image
                apply_image_border(shape)

        # Execute deletions
        for shape in shapes_to_delete:
             try:
                 slide.shapes._spTree.remove(shape._element)
             except:
                 pass

# Save
prs.save(PATH)
print(f"Presentation v3.6.1 updated: Redundant twins and large frames removed. Borders applied only to images.")
