from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.1.pptx"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: PATH")
    exit(1)

prs = Presentation(PATH)

def apply_luxury_border(shape, color=(222, 255, 154), width_pt=3):
    line = shape.line
    line.color.rgb = RGBColor(*color)
    line.width = Pt(width_pt)

# 1. Slide 4 (Index 3): Project Overview
if len(prs.slides) > 3:
    slide = prs.slides[3]
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            # Resize and Center
            shape.width = Inches(8.5)
            shape.left = Inches(0.75)
            shape.top = Inches(1.2)
            # Add luxury border
            apply_luxury_border(shape)
        elif hasattr(shape, "text") and "전통적 우드" in shape.text:
            # Move text down a bit
            shape.top = Inches(6.8)
            shape.left = Inches(0.75)
            shape.width = Inches(8.5)

# 2. Slide 22 (Index 21): Strategic Risk
if len(prs.slides) > 21:
    slide = prs.slides[21]
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            shape.width = Inches(8.5)
            shape.left = Inches(0.75)
            shape.top = Inches(1.2)
            apply_luxury_border(shape)
        elif hasattr(shape, "text") and "긍정적 지표" in shape.text:
            shape.top = Inches(6.8)
            shape.left = Inches(0.75)
            shape.width = Inches(8.5)

# 3. Slide 23 (Index 22): Differentiation Strategy
if len(prs.slides) > 22:
    slide = prs.slides[22]
    # Remove borders/frames
    for shape in list(slide.shapes):
        if shape.shape_type == 1: # Rectangle (likely the frame)
             # But wait, we only want to remove the ones that look like frames
             if shape.width > Inches(8): # Large rectangle
                 slide.shapes._spTree.remove(shape._element)
        elif shape.shape_type == 13: # Picture
            # Remove border from picture if any
            shape.line.fill.background() 
            # Center it
            shape.width = Inches(7.5)
            shape.left = Inches(1.25)
            shape.top = Inches(1.5)
        elif hasattr(shape, "text") and "유동인구" in shape.text:
            shape.top = Inches(6.8)

# Save
prs.save(PATH)
print(f"Presentation v3.6.1 updated with luxury borders and layout fixes.")
