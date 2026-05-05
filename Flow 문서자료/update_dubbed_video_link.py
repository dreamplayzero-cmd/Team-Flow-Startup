from pptx import Presentation
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
NEW_DUBBED_VIDEO = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\프로젝트+시연영상.mp4"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

# --- Update Video Link (Slide 19 only) ---
# Slide 19 is index 18
if len(prs.slides) > 18:
    slide = prs.slides[18]
    
    # 1. Find the link text and update its hyperlink
    target_found = False
    for shape in slide.shapes:
        if hasattr(shape, "text") and "시연영상" in shape.text:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    # Update Hyperlink to the new dubbed version
                    video_url = f"file:///{NEW_DUBBED_VIDEO.replace(chr(92), '/')}"
                    run.hyperlink.address = video_url
                    target_found = True

    if not target_found:
        print("Link text not found on Slide 19.")
else:
    print(f"Slide 19 not found. Total slides: {len(prs.slides)}")

# Save
prs.save(PATH)
print(f"Presentation v3.6 Updated with Dubbed Video Link saved to: {PATH}")
