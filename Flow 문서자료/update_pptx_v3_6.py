from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.5.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
DEMO_VIDEO_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\모바일 시연영상.mp4"

# --- Image Paths ---
KOREAN_LOGIC_IMG = os.path.join(IMG_BASE, "korean_score_logic_v3_6_1777365179186.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def set_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_body_text(slide, text, left, top, width, height, size=14, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    return txBox

# 1. Update Slide 12 (Harness Engineering)
# Slide index 11
if len(prs.slides) > 11:
    slide = prs.slides[11]
    add_body_text(slide, 
        "하네스(Harness) 엔지니어링 개요:\n"
        "- 비정형 데이터 정제 및 Gemini API 최적화 전처리 엔진\n"
        "- 가중치 산정 로직(공공 6:4 SNS) 실시간 적용\n"
        "- 데이터 파이프라인의 오케스트레이션 제어", 
        Inches(0.5), Inches(5.5), Inches(8), Inches(1.5), size=14, color=(222, 255, 154))

# 2. Update Risk Analysis & Core Logic Slides (Translate Titles)
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "STRATEGIC RISK & INSIGHT ANALYSIS" in shape.text:
                shape.text = "전략적 리스크 및 인사이트 분석"
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
            if "POSITIVE INDICATORS" in shape.text:
                shape.text = "긍정적 지표"
            if "RISK FACTORS" in shape.text:
                shape.text = "리스크 요인"
            if "THE INSIGHT ENGINE: CORE LOGIC" in shape.text:
                shape.text = "인사이트 엔진: 핵심 로직"
                # Replace image with Korean version if found
                for s in slide.shapes:
                    if s.shape_type == 13: # Picture
                        # Replace picture
                        left, top, width, height = s.left, s.top, s.width, s.height
                        slide.shapes._spTree.remove(s._element)
                        if os.path.exists(KOREAN_LOGIC_IMG):
                            slide.shapes.add_picture(KOREAN_LOGIC_IMG, left, top, width=width)

# 3. Slide 18/19 Demo Video Split
# Find slide 18 (index 17)
if len(prs.slides) > 17:
    slide18 = prs.slides[17]
    # Remove Demo Video Link from Slide 18 if it exists (Keep Promotion)
    for shape in list(slide18.shapes):
        if hasattr(shape, "text") and ("모바일 시연영상" in shape.text or "PLAY DEMO VIDEO" in shape.text):
            slide18.shapes._spTree.remove(shape._element)

    # Create Slide 19 (New) for Demo Video
    slide19 = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_theme(slide19)
    title = slide19.shapes.title
    title.text = "서비스 시연 영상 (모바일)"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
    
    # Move slide 19 to after slide 18
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(18, sldId)
    
    # Add Video Link to Slide 19
    video_url = f"file:///{DEMO_VIDEO_PATH.replace(chr(92), '/')}"
    add_body_text(slide19, "▶ 모바일 시연영상 재생 (클릭)", Inches(1), Inches(2), Inches(8), Inches(1), size=24, color=(222, 255, 154))
    btn = slide19.shapes.add_shape(1, Inches(1), Inches(3), Inches(4), Inches(1))
    btn.text = "시연영상 시작"
    btn.click_action.hyperlink.address = video_url

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.6 saved to: {OUTPUT_PATH}")
