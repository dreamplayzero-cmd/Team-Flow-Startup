from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v2.7.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.0.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
DNA_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\backend\assets\visual_dna"

# --- Image Mapping ---
WEB_IMG = os.path.join(IMG_BASE, "web_frontend_v3_premium_1777356658037.png")
MOBILE_IMG = os.path.join(IMG_BASE, "mobile_frontend_v3_premium_1777356694189.png")
SCORE_IMG = os.path.join(IMG_BASE, "score_logic_v3_premium_1777356820324.png")
CLUSTER_IMG = os.path.join(IMG_BASE, "clustering_visualization_premium_1777346047876.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

# Helper: Dark theme application
def apply_dark_theme(slide):
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 10) # Very dark gray
    
    # Update text colors
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(245, 245, 245) # White

# Apply theme to all existing slides
for slide in prs.slides:
    apply_dark_theme(slide)

# --- Add New Content Slides ---

def add_new_slide(title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    apply_dark_theme(slide)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154) # Gold
    return slide

# 1. Web Frontend Capture
slide = add_new_slide("WEB FRONTEND: DISTRICT ANALYSIS DASHBOARD")
if os.path.exists(WEB_IMG):
    slide.shapes.add_picture(WEB_IMG, Inches(1), Inches(1.5), width=Inches(8))

# 2. Mobile Frontend Capture
slide = add_new_slide("MOBILE FRONTEND: INSIGHT & SCOUTING")
if os.path.exists(MOBILE_IMG):
    slide.shapes.add_picture(MOBILE_IMG, Inches(2), Inches(1.2), height=Inches(5.5))

# 3. Score Summation Logic
slide = add_new_slide("INTEGRATED SCORING ENGINE (v2.7)")
if os.path.exists(SCORE_IMG):
    slide.shapes.add_picture(SCORE_IMG, Inches(1), Inches(1.5), width=Inches(8))

# 4. Recommended Store Images (Visual DNA)
slide = add_new_slide("RECOMMENDED STORE DESIGN (VISUAL DNA)")
dna_files = [
    "SS_ID_02_Seongsu_Cafe_IndustrialVintage.jpg",
    "HN_MC_01_Hannam_Dining_ModernChic.jpg",
    "SS_WW_01_Seongsu_Cafe_WarmWood.jpg"
]
for i, f in enumerate(dna_files):
    path = os.path.join(DNA_PATH, f)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(0.5 + i*3.1), Inches(2), width=Inches(3))

# 5. Clustering Table
slide = add_new_slide("AI CLUSTERING & FUTURE SCENARIOS")
if os.path.exists(CLUSTER_IMG):
    slide.shapes.add_picture(CLUSTER_IMG, Inches(1), Inches(1.5), width=Inches(8))

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.0 saved to: {OUTPUT_PATH}")
