from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.3.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.5.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
VIDEO_PATH = r"c:\Users\User\Desktop\Flow 발표자료\성수의_부상과_군림.mp4"

# --- Image Paths ---
IMG_ARCH = os.path.join(IMG_BASE, "harness_logic_architecture_v3_4_1777363383609.png")
IMG_MOBILE = os.path.join(IMG_BASE, "mobile_gis_mockup_v3_4_1777363524637.png")
IMG_WEB = os.path.join(IMG_BASE, "web_depth_analysis_v3_4_1777363711541.png")
IMG_TEAM = os.path.join(IMG_BASE, "team_synergy_v3_4_1777363878514.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def set_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_text(slide, text, color=(222, 255, 154)):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)

def add_body_text(slide, text, left, top, width, height, size=18, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    return txBox

# --- Re-create Executive Summary at the START ---

# Slide 1: Executive Title
# In python-pptx, you can't easily insert at index 0, so we append and move (complex)
# or we just recreate the WHOLE deck starting with new slides.
# Actually, we can move slides by manipulating the sldIdLst.

# Let's just create a new presentation, add summary, THEN copy v3.3 slides? 
# No, copying slides is hard.
# I'll add the summary slides at the END and the user can move them, 
# OR I'll try to insert them at the beginning using slide ID manipulation.

def insert_slide_at(prs, layout, index):
    slide = prs.slides.add_slide(layout)
    # Move slide to index
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(index, sldId)
    return slide

# Insert 6 Executive Slides at index 1 to 6 (after original title) or just 0 to 5.
# Let's start at index 1.

# 1. Architecture (Harness Logic)
slide = insert_slide_at(prs, prs.slide_layouts[5], 1)
set_dark_theme(slide)
add_title_text(slide, "EXECUTIVE SUMMARY: HARNESS LOGIC")
add_body_text(slide, "Bucket -> Harness -> Gemini\n단순 도식을 넘어선 엔지니어링 아키텍처 구현.", Inches(0.5), Inches(1.2), Inches(9), Inches(1))
if os.path.exists(IMG_ARCH):
    slide.shapes.add_picture(IMG_ARCH, Inches(1), Inches(2.5), width=Inches(8))

# 2. The Face (Mobile)
slide = insert_slide_at(prs, prs.slide_layouts[5], 2)
set_dark_theme(slide)
add_title_text(slide, "THE FACE: MOBILE STRATEGY")
add_body_text(slide, "9,900 KRW Strategy\n현장 중심 UI 및 시연 영상 연동.", Inches(0.5), Inches(1.2), Inches(5), Inches(1))
if os.path.exists(IMG_MOBILE):
    slide.shapes.add_picture(IMG_MOBILE, Inches(5.5), Inches(1.5), height=Inches(5))
if os.path.exists(VIDEO_PATH):
    video_url = f"file:///{VIDEO_PATH.replace(chr(92), '/')}"
    btn = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(3), Inches(0.8))
    btn.text = "PLAY DEMO VIDEO"
    btn.click_action.hyperlink.address = video_url

# 3. The Heart (Web)
slide = insert_slide_at(prs, prs.slide_layouts[5], 3)
set_dark_theme(slide)
add_title_text(slide, "THE HEART: WEB ANALYSIS")
add_body_text(slide, "29,000 KRW Strategy\n정밀 데이터 시트 및 심층 분석.", Inches(5), Inches(2), Inches(4.5), Inches(2))
if os.path.exists(IMG_WEB):
    slide.shapes.add_picture(IMG_WEB, Inches(0.5), Inches(1.5), width=Inches(4.2))

# 4. Synergy
slide = insert_slide_at(prs, prs.slide_layouts[5], 4)
set_dark_theme(slide)
add_title_text(slide, "BUSINESS IMPACT: SYNERGY")
add_body_text(slide, "Double Strategy Synergy\nTeam Flow Integrated Model.", Inches(1), Inches(1.5), Inches(8), Inches(1))
if os.path.exists(IMG_TEAM):
    slide.shapes.add_picture(IMG_TEAM, Inches(1), Inches(2.5), width=Inches(8))

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.5 saved to: {OUTPUT_PATH}")
