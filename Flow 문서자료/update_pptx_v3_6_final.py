from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

def add_styled_body_text(slide, text, left, top, width, height, size=14, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    return txBox

def set_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

# --- Targeted Updates (6, 7, 19, 21) ---

# Slide 6 (index 5): Commercial Activity Index
if len(prs.slides) > 5:
    slide = prs.slides[5]
    add_styled_body_text(slide, 
        "핵심 분석 요약:\n"
        "- 성수(Vintage Wood): 주말 및 야간 유동인구 집중형 상권 (Experience-Driven)\n"
        "- 한남(Minimal Modern): 주중 일정한 활동 지수를 유지하는 안정형 상권 (Value-Driven)\n"
        "- 시사점: 타겟 고객의 방문 목적에 따른 최적 운영 시간대 도출 가능",
        Inches(1), Inches(6.8), Inches(8), Inches(1.5), size=12, color=(222, 255, 154))

# Slide 7 (index 6): Visual DNA Matching
if len(prs.slides) > 6:
    slide = prs.slides[6]
    add_styled_body_text(slide, 
        "시각적 분석 결과:\n"
        "- 성수: #팝업스토어, #빈티지, #공장개조 등 '변주와 파격'의 시각 정체성\n"
        "- 한남: #미니멀, #갤러리, #정제된감성 등 '절제와 조화'의 시각 정체성\n"
        "- 결과: 각 구역의 시각적 기대치에 부합하는 인테리어 DNA 매칭 (Match Score 반영)",
        Inches(1), Inches(6.8), Inches(8), Inches(1.5), size=12, color=(222, 255, 154))

# Slide 19 (index 18): Service Demo Video (Make it sensible/premium)
if len(prs.slides) > 18:
    slide = prs.slides[18]
    # Redesign: Add a dark marble frame or better layout
    set_dark_theme(slide)
    # Clear and recreate title
    for shape in list(slide.shapes):
        if shape.shape_type == 14: # Title
             shape.text = "THE FACE: SERVICE DEMONSTRATION"
             shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
    
    add_styled_body_text(slide, 
        "MOBILE GIS & SCOUTING INTERFACE\n"
        "현장의 감각을 데이터로 전환하는 직관적 UX 시연", 
        Inches(0.5), Inches(1.5), Inches(5), Inches(1), size=18, color=(245, 245, 245))
    
    # Ensure the video link button is styled nicely
    for shape in slide.shapes:
        if hasattr(shape, "text") and "시연영상" in shape.text:
             if shape.shape_type == 1: # Rectangle/Button
                 shape.fill.solid()
                 shape.fill.fore_color.rgb = RGBColor(33, 65, 85) # Slate Gray
                 shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)

# Slide 21 (index 20): Dream Team
if len(prs.slides) > 20:
    slide = prs.slides[20]
    add_styled_body_text(slide, 
        "프로젝트 전문성 보강:\n"
        "- Member A: 엔지니어링 아키텍처 및 데이터 무결성 보장\n"
        "- Member B: 상권 시각 분석 알고리즘 고도화 및 브랜딩 총괄\n"
        "- 협업 시너지: 기술적 Harness와 감각적 Face의 완벽한 결합",
        Inches(1), Inches(6.5), Inches(8), Inches(1.5), size=12, color=(222, 255, 154))

# Save
prs.save(PATH)
print(f"Presentation v3.6 Updated (Slides 6, 7, 19, 21) saved to: {PATH}")
