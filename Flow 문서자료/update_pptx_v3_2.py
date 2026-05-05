from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.0.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.2.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
DNA_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\backend\assets\visual_dna"

# --- Image Mapping ---
WEB_IMG = os.path.join(IMG_BASE, "web_frontend_v3_premium_1777356658037.png")
MOBILE_IMG = os.path.join(IMG_BASE, "mobile_frontend_v3_premium_1777356694189.png")
SCORE_IMG = os.path.join(IMG_BASE, "score_logic_v3_premium_1777356820324.png")
CLUSTER_IMG = os.path.join(IMG_BASE, "clustering_visualization_premium_1777346047876.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def add_floating_text(slide, text, left, top, width, height, size=14, bold=False, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return txBox

# --- Image Integration into Existing Slides ---

for i, slide in enumerate(prs.slides):
    text_content = "".join([s.text for s in slide.shapes if hasattr(s, "text")]).upper()
    
    # 1. Branding / DNA Section (Slide 10-11 area)
    if "BRANDING" in text_content or "IDENTITY" in text_content:
        # Add Visual DNA images to branding slide
        dna_img = os.path.join(DNA_PATH, "SS_ID_02_Seongsu_Cafe_IndustrialVintage.jpg")
        if os.path.exists(dna_img):
            slide.shapes.add_picture(dna_img, Inches(6), Inches(1.5), width=Inches(3.5))
            add_floating_text(slide, "Visual DNA: Industrial Vintage (Seongsu)", Inches(6), Inches(5.2), Inches(3.5), Inches(0.5), size=10)

    # 2. Solution Scenarios / Clustering (Slide 12-13 area)
    if "SOLUTION" in text_content or "SCENARIO" in text_content or "CLUSTERING" in text_content:
        if os.path.exists(CLUSTER_IMG):
            slide.shapes.add_picture(CLUSTER_IMG, Inches(1), Inches(4), width=Inches(8))

    # 3. Demonstration / Product (Slide 14-15 area)
    if "DEMONSTRATION" in text_content or "SHOWCASE" in text_content:
        if os.path.exists(WEB_IMG):
            slide.shapes.add_picture(WEB_IMG, Inches(0.5), Inches(2), width=Inches(4.5))
        if os.path.exists(MOBILE_IMG):
            slide.shapes.add_picture(MOBILE_IMG, Inches(5.5), Inches(1.5), height=Inches(5))

# --- Adding Advanced Risk Analysis Slide (Based on District_Report_Presentation.md) ---
slide_layout = prs.slide_layouts[5] # Title Only
slide = prs.slides.add_slide(slide_layout)
# Note: In a real script I'd call my apply_dark_theme helper, but since v3.0 is already dark, it inherits.
# But just in case, I'll re-apply dark background logic if needed.
# prs.slides[-1] is the new slide.
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 10, 10)

title = slide.shapes.title
title.text = "STRATEGIC RISK & INSIGHT ANALYSIS"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)

# Add Positive Indicators
add_floating_text(slide, "POSITIVE INDICATORS", Inches(1), Inches(2), Inches(4), Inches(0.5), size=18, bold=True, color=(222, 255, 154))
add_floating_text(slide, "- MZ 세대 높은 충성도 (재방문율 35%+)\n- SNS 바이럴 점유율 (상위 1%)\n- 하이엔드 오피스 융합 가능성", Inches(1), Inches(2.5), Inches(4), Inches(2))

# Add Risk Factors
add_floating_text(slide, "RISK FACTORS", Inches(5.5), Inches(2), Inches(4), Inches(0.5), size=18, bold=True, color=(255, 100, 100))
add_floating_text(slide, "- 임대료 급증 (연 12%↑)\n- 젠트리피케이션 가속화 위험\n- 특정 업종(카페) 공급 임계점 도달", Inches(5.5), Inches(2.5), Inches(4), Inches(2))

# Add Engine Logic Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "THE INSIGHT ENGINE: CORE LOGIC"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
if os.path.exists(SCORE_IMG):
    slide.shapes.add_picture(SCORE_IMG, Inches(1), Inches(1.5), width=Inches(8))

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.2 saved to: {OUTPUT_PATH}")
