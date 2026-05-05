from pptx import Presentation
from pptx.util import Inches
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.2.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.3.pptx"

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

# 1. Identify and Remove Redundant Slides
# In v3.2, I added 2 new slides at the end (Risk Analysis and Engine Logic).
# In v3.0 (source for v3.2), I added 5 slides at the end.
# So slides -7 to -3 (approx) are redundant duplicates of integrated content.
# Let's count them.
total_slides = len(prs.slides)
print(f"Initial Slide Count: {total_slides}")

# I want to keep the "Strategic Risk" and "Engine Logic" slides (the last 2).
# I want to remove the 5 slides added in v3.0 which are now redundant.
# These would be indices total_slides - 7 to total_slides - 3.

slides_to_remove = []
# Identify slides that were just "Demo" slides (Title Only with one image)
for i in range(total_slides):
    slide = prs.slides[i]
    text_content = "".join([s.text for s in slide.shapes if hasattr(s, "text")]).upper()
    
    # Redundant "Demonstration" slides usually have very specific titles
    if any(k in text_content for k in ["WEB FRONTEND: DISTRICT", "MOBILE FRONTEND: INSIGHT", "INTEGRATED SCORING ENGINE (v2.7)", "RECOMMENDED STORE DESIGN", "AI CLUSTERING & FUTURE SCENARIOS"]):
        # But wait, I added these in v3.0. If they are exactly these titles, they are the "extra" slides.
        slides_to_remove.append(i)

# Remove slides in reverse order
for i in sorted(slides_to_remove, reverse=True):
    # prs.slides._sldIdLst.remove(prs.slides[i]._element) # This is the hacky way to remove
    # Actually, a better way:
    rId = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[i]

print(f"Removed {len(slides_to_remove)} redundant slides.")

# 2. Fix Overlaps on remaining slides
for slide in prs.slides:
    # If a slide has too many pictures, try to reorganize
    pics = [s for s in slide.shapes if s.shape_type == 13] # 13 is Picture
    if len(pics) >= 2:
        # Check for overlap
        for i in range(len(pics)):
            for j in range(i + 1, len(pics)):
                p1 = pics[i]
                p2 = pics[j]
                # Simple overlap check
                if abs(p1.left - p2.left) < Inches(1) and abs(p1.top - p2.top) < Inches(1):
                    # Move p2 to the right
                    p2.left = p1.left + p1.width + Inches(0.2)

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.3 saved to: {OUTPUT_PATH}")
