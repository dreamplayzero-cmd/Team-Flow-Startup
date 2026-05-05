from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.3.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.5.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
NEW_VIDEO_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\모바일 시연영상.mp4"

# --- Image Paths ---
IMG_ARCH = os.path.join(IMG_BASE, "harness_logic_architecture_v3_4_1777363383609.png")
IMG_MOBILE = os.path.join(IMG_BASE, "mobile_gis_mockup_v3_4_1777363524637.png")
IMG_WEB = os.path.join(IMG_BASE, "web_depth_analysis_v3_4_1777363711541.png")
IMG_TEAM = os.path.join(IMG_BASE, "team_synergy_v3_4_1777363878514.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def set_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_text(slide, text, color=(222, 255, 154)):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)

def add_body_text(slide, text, left, top, width, height, size=18, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    return txBox

def insert_slide_at(prs, layout, index):
    slide = prs.slides.add_slide(layout)
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(index, sldId)
    return slide

# 1. Insert Executive Summary AFTER TOC (index 2 onwards)
# Slide 0: Title, Slide 1: TOC
# Start inserting from index 2.

# Executive 1: Architecture
slide = insert_slide_at(prs, prs.slide_layouts[5], 2)
set_dark_theme(slide)
add_title_text(slide, "EXECUTIVE SUMMARY: HARNESS LOGIC")
if os.path.exists(IMG_ARCH):
    slide.shapes.add_picture(IMG_ARCH, Inches(1), Inches(1.5), width=Inches(8))

# Executive 2: Mobile Strategy
slide = insert_slide_at(prs, prs.slide_layouts[5], 3)
set_dark_theme(slide)
add_title_text(slide, "THE FACE: MOBILE SCOUTING")
if os.path.exists(IMG_MOBILE):
    slide.shapes.add_picture(IMG_MOBILE, Inches(5.5), Inches(1.2), height=Inches(5.5))
add_body_text(slide, "9,900 KRW Strategy\n현장 중심 UI 및 시연 연동.", Inches(0.5), Inches(1.5), Inches(4.5), Inches(2))

# Executive 3: Web Strategy
slide = insert_slide_at(prs, prs.slide_layouts[5], 4)
set_dark_theme(slide)
add_title_text(slide, "THE HEART: WEB ANALYSIS")
if os.path.exists(IMG_WEB):
    slide.shapes.add_picture(IMG_WEB, Inches(0.5), Inches(1.5), width=Inches(5))
add_body_text(slide, "29,000 KRW Strategy\n심층 데이터 리포트.", Inches(6), Inches(2), Inches(3.5), Inches(2))

# 2. Update Video Link on Demonstration Slide
target_text = "DEMONSTRATION"
video_url = f"file:///{NEW_VIDEO_PATH.replace(chr(92), '/')}"

for slide in prs.slides:
    text_content = "".join([s.text for s in slide.shapes if hasattr(s, "text")]).upper()
    if target_text in text_content:
        # Found it! Add or update the link
        add_body_text(slide, "▶ 시연 영상 재생 (모바일 시연영상.mp4)", Inches(0.5), Inches(5.5), Inches(8), Inches(1), size=14, color=(222, 255, 154))
        btn = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(3), Inches(0.8))
        btn.text = "모바일 시연영상 재생"
        btn.click_action.hyperlink.address = video_url

# 3. Final Theme Check for TOC
# Ensure slide 1 (TOC) is dark
if len(prs.slides) > 1:
    set_dark_theme(prs.slides[1])
    for shape in prs.slides[1].shapes:
        if hasattr(shape, "text_frame"):
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor(245, 245, 245)

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.5 Fixed saved to: {OUTPUT_PATH}")
