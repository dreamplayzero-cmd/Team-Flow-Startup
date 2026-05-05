from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Configuration ---
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.4_Executive.pptx"
IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
VIDEO_PATH = r"c:\Users\User\Desktop\Flow 발표자료\성수의_부상과_군림.mp4"

# --- Image Paths ---
IMG_ARCH = os.path.join(IMG_BASE, "harness_logic_architecture_v3_4_1777363383609.png")
IMG_MOBILE = os.path.join(IMG_BASE, "mobile_gis_mockup_v3_4_1777363524637.png")
IMG_WEB = os.path.join(IMG_BASE, "web_depth_analysis_v3_4_1777363711541.png")
IMG_TEAM = os.path.join(IMG_BASE, "team_synergy_v3_4_1777363878514.png")

# --- Initialize ---
prs = Presentation()

def set_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_text(slide, text, color=(222, 255, 154)): # Sovereign Gold
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)

def add_body_text(slide, text, left, top, width, height, size=18, color=(245, 245, 245)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    return txBox

# 1. Slide: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_dark_theme(slide)
title = slide.shapes.title
title.text = "Youth Startup Flow - Final Report"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
subtitle = slide.placeholders[1]
subtitle.text = "Sovereign Insight: Strategic Multi-Platform Ecosystem"
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(245, 245, 245)

# 2. Slide: Architecture (Harness Logic)
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_dark_theme(slide)
add_title_text(slide, "ARCHITECTURE: HARNESS LOGIC")
add_body_text(slide, "Bucket -> Harness -> Gemini\n단순 도식을 넘어선 엔지니어링 아키텍처 구현.", Inches(0.5), Inches(1.2), Inches(9), Inches(1))
if os.path.exists(IMG_ARCH):
    slide.shapes.add_picture(IMG_ARCH, Inches(1), Inches(2.5), width=Inches(8))

# 3. Slide: The Face (Mobile)
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_dark_theme(slide)
add_title_text(slide, "THE FACE: MOBILE DEMO")
add_body_text(slide, "9,900 KRW Strategy\n시연 영상과 시각적 궤를 같이하는 현장 중심 UI 배치.", Inches(0.5), Inches(1.2), Inches(5), Inches(1))
if os.path.exists(IMG_MOBILE):
    slide.shapes.add_picture(IMG_MOBILE, Inches(5.5), Inches(1.5), height=Inches(5))

# Add Video Link or Placeholder
if os.path.exists(VIDEO_PATH):
    video_url = f"file:///{VIDEO_PATH.replace(chr(92), '/')}"
    add_body_text(slide, f"▶ Click to Play Demo Video\n({os.path.basename(VIDEO_PATH)})", Inches(0.5), Inches(5.5), Inches(4), Inches(1), size=14, color=(222, 255, 154))
    btn = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(3), Inches(0.8)) # Rectangle
    btn.text = "PLAY DEMO VIDEO"
    btn.click_action.hyperlink.address = video_url

# 4. Slide: The Heart (Web)
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_dark_theme(slide)
add_title_text(slide, "THE HEART: WEB DEPTH ANALYSIS")
add_body_text(slide, "29,000 KRW Strategy\n29,000원의 가치를 증명하는 정밀 데이터 시트 구축.", Inches(5), Inches(2), Inches(4.5), Inches(2))
if os.path.exists(IMG_WEB):
    slide.shapes.add_picture(IMG_WEB, Inches(0.5), Inches(1.5), width=Inches(4.2))

# 5. Slide: Business Impact
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_dark_theme(slide)
add_title_text(slide, "BUSINESS IMPACT: SYNERGY")
add_body_text(slide, "Double Strategy Synergy\nMobile(Face) + Web(Heart) 통합 비즈니스 모델.", Inches(1), Inches(1.5), Inches(8), Inches(1))
if os.path.exists(IMG_TEAM):
    slide.shapes.add_picture(IMG_TEAM, Inches(1), Inches(2.5), width=Inches(8))

# 6. Slide: Closing
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_dark_theme(slide)
title = slide.shapes.title
title.text = "Q&A & CONTACT"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(222, 255, 154)
subtitle = slide.placeholders[1]
subtitle.text = "Thank you for your attention.\nTeam Flow: Sovereignty in Insight"
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(245, 245, 245)

# Save
prs.save(OUTPUT_PATH)
print(f"Executive Presentation v3.4 saved to: {OUTPUT_PATH}")
