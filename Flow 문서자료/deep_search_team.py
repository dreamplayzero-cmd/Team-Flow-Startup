import os
from pptx import Presentation

search_path = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료"
results = []

for filename in os.listdir(search_path):
    if filename.endswith(".pptx"):
        full_path = os.path.join(search_path, filename)
        try:
            prs = Presentation(full_path)
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if "팀 리더" in shape.text or "Team Leader" in shape.text or "Member C" in shape.text:
                            results.append({
                                "file": filename,
                                "slide": i + 1,
                                "text": shape.text
                            })
        except Exception as e:
            print(f"Error reading {filename}: {e}")

for res in results:
    print(f"--- Found in {res['file']} (Slide {res['slide']}) ---")
    print(res['text'])
    print("-" * 30)
