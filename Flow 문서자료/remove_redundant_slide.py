from pptx import Presentation
import os

PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.3.pptx"

if os.path.exists(PATH):
    prs = Presentation(PATH)
    # Slide 18 is index 17. Slide 20 is index 19.
    # After deleting index 17, the original index 19 becomes index 18.
    
    # I'll delete Slide 18 (index 17) which has the "(v2.7)" tag.
    # This keeps Slide 20 ("THE INSIGHT ENGINE: CORE LOGIC") as the main engine logic slide.
    
    idx_to_remove = 17 
    rId = prs.slides._sldIdLst[idx_to_remove].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx_to_remove]
    
    prs.save(PATH)
    print(f"Removed redundant Slide 18. Saved to {PATH}")
else:
    print("File not found")
