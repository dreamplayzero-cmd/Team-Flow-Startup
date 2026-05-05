from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# --- Configuration ---
SOURCE_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"
OUTPUT_PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.1.pptx"

IMG_BASE = r"C:\Users\User\.gemini\antigravity\brain\aa691bd8-c4b7-4244-8015-3b2e0b6fdca5"
IMG_SLIDE_4 = os.path.join(IMG_BASE, "media__1777422537006.png")
IMG_SLIDE_22 = os.path.join(IMG_BASE, "media__1777422634373.png")
IMG_SLIDE_23 = os.path.join(IMG_BASE, "media__1777422650754.png")

# --- Initialize ---
if not os.path.exists(SOURCE_PATH):
    print(f"Source not found: {SOURCE_PATH}")
    exit(1)

prs = Presentation(SOURCE_PATH)

def clear_slide_content(slide):
    # Keep only the Title or large Title placeholders
    shapes_to_keep = []
    for shape in slide.shapes:
        is_title = False
        if shape.name.startswith("Title"):
            is_title = True
        elif shape.is_placeholder:
            # Placeholder type 1 is TITLE
            try:
                if shape.placeholder_format.type == 1:
                    is_title = True
            except:
                pass
        
        if is_title:
            shapes_to_keep.append(shape)
    
    # Remove everything else
    for shape in list(slide.shapes):
        if shape not in shapes_to_keep:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass

def add_image_with_border_and_text(slide, img_path, description):
    # 1. Add Image
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(6.5)
    pic = slide.shapes.add_picture(img_path, left, top, width=width)
    
    # 2. Add Green Border to Image
    line = pic.line
    line.color.rgb = RGBColor(222, 255, 154) # Sovereign Gold / Neon Green
    line.width = Pt(3)
    
    # 3. Add Description Text to the side
    tx_left = Inches(7.2)
    tx_top = Inches(2.2)
    tx_width = Inches(2.5)
    tx_height = Inches(3)
    txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(245, 245, 245)
    p.font.name = "Noto Sans KR"
    p.font.bold = True

# 1. Slide 4 (Project Overview)
if len(prs.slides) > 3:
    slide = prs.slides[3]
    clear_slide_content(slide)
    add_image_with_border_and_text(slide, IMG_SLIDE_4, "성수·한남동의 감성과 데이터를 결합한 AI 상권 분석 플랫폼 리포트")

# 2. Slide 22 (Strategic Risk)
if len(prs.slides) > 21:
    slide = prs.slides[21]
    clear_slide_content(slide)
    add_image_with_border_and_text(slide, IMG_SLIDE_22, "리스크 요인과 긍정적 지표를 종합한 상권 군집 분석 및 최종 결론")

# 3. Slide 23 (Differentiation Strategy)
if len(prs.slides) > 22:
    slide = prs.slides[22]
    clear_slide_content(slide)
    add_image_with_border_and_text(slide, IMG_SLIDE_23, "유동인구, 트렌드, 수요, 경쟁의 4대 지표를 통합한 핵심 인사이트 엔진")

# Save
prs.save(OUTPUT_PATH)
print(f"Presentation v3.6.1 updated successfully from v3.6 base.")
