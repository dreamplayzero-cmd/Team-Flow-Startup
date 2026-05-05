from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Configuration ---
PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.6.pptx"

# --- Initialize ---
if not os.path.exists(PATH):
    print(f"Source not found: {PATH}")
    exit(1)

prs = Presentation(PATH)

# --- Targeted Update (Slide 19 only) ---
# Slide 19 is index 18
if len(prs.slides) > 18:
    slide = prs.slides[18]
    
    # Sovereign Gold / Neon Green color
    NEON_GREEN = RGBColor(222, 255, 154) # #DEFF9A
    
    # 1. Find the mobile mockup image to get its position
    mockup_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            mockup_shape = shape
            break
            
    if mockup_shape:
        # 2. Add the "Frame" behind it
        # Calculate frame size (slightly larger than mockup)
        padding = Inches(0.15)
        left = mockup_shape.left - padding
        top = mockup_shape.top - padding
        width = mockup_shape.width + (padding * 2)
        height = mockup_shape.height + (padding * 2)
        
        # Add rectangle
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        frame.fill.background() # No fill
        frame.line.color.rgb = NEON_GREEN
        frame.line.width = Pt(3.5)
        
        # 3. Send frame to back (so it's behind the mockup)
        # In python-pptx, new shapes are added to the front. We need to move it back.
        # Simplest way is to remove and re-add other shapes, or use the z-order logic if available.
        # Actually, let's just make sure the mockup is brought to front.
        slide.shapes._spTree.remove(mockup_shape._element)
        slide.shapes._spTree.append(mockup_shape._element)

# Save
prs.save(PATH)
print(f"Presentation v3.6 Enhanced (Slide 19 Frame) saved to: {PATH}")
