from pptx import Presentation
import os

PATH = r"c:\Users\User\Desktop\Team flow -Youth Startup Flow\Flow 문서자료\Flow_Final_Project_Report_v3.3.pptx"

if os.path.exists(PATH):
    prs = Presentation(PATH)
    # Check slide 18 (index 17) and slide 20 (index 19)
    for idx in [17, 19]:
        if idx < len(prs.slides):
            slide = prs.slides[idx]
            print(f"Slide {idx+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"  - Text: {shape.text[:100]}...")
                if shape.shape_type == 13: # Picture
                    print(f"  - [Picture Found]")
        else:
            print(f"Slide {idx+1} does not exist.")
else:
    print("File not found")
